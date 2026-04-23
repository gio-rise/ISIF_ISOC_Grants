"""
Build editable PPTX for the ISIF Capability concept video.
Uses python-pptx for native editable text/shapes (vs MARP's rasterized output).
Applies RISE Style Guide 2019 brand system.

Run from project root:
    uv run docs/presentations/isif-capability-video/build_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# --- RISE Style Guide 2019: brand colors ---
RISE_PINK   = RGBColor(0xF0, 0x4F, 0x6C)
RISE_BLACK  = RGBColor(0x33, 0x32, 0x32)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BG_NAVY     = RGBColor(0x10, 0x4B, 0x84)
BG_DARK     = RGBColor(0x0F, 0x10, 0x29)
TEAL        = RGBColor(0x27, 0xB9, 0x9A)
# Utility greys (not brand colors; used only for borders and muted text)
GREY_700    = RGBColor(0x4A, 0x55, 0x68)
GREY_500    = RGBColor(0x71, 0x80, 0x96)
GREY_300    = RGBColor(0xCB, 0xD5, 0xE0)
GREY_100    = RGBColor(0xF7, 0xFA, 0xFC)

# --- Typography (RISE brand: Maison Neue; universal fallback: Calibri) ---
FONT_HEAD = "Maison Neue Demi"
FONT_BODY = "Maison Neue Book"
FONT_FALLBACK = "Calibri"  # PowerPoint will use this if Maison Neue is absent

# --- Slide geometry: 16:9 widescreen, 13.333 x 7.5 inches ---
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

OUTPUT = Path(__file__).parent / "slides.pptx"


# ---------- helpers ----------

def set_slide_bg(slide, color):
    """Apply solid background color to a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, *,
                font=FONT_BODY, size=14, bold=False, italic=False,
                color=RISE_BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                line_spacing=1.15):
    """Add a text box with styled text. Returns the shape."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing

    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill_color, *,
             line_color=None, line_weight=0):
    """Add a filled rectangle. Optionally with border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is not None and line_weight > 0:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_weight)
    else:
        shape.line.fill.background()  # no border
    shape.shadow.inherit = False
    return shape


def add_top_border(slide, left, top, width, color, thickness_pt=3):
    """Add a thin accent bar on top of a region."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                  width, Pt(thickness_pt))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_left_border(slide, left, top, height, color, thickness_pt=4):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top,
                                  Pt(thickness_pt), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def add_rise_mark(slide, color=GREY_500):
    """Small RISE wordmark bottom-left. Placeholder until logo image provided."""
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(1.5), Inches(0.3),
                "RISE", font=FONT_HEAD, size=10, bold=True,
                color=color, align=PP_ALIGN.LEFT)


def add_source(slide, text, color=GREY_500):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
                text, font=FONT_BODY, size=9, italic=True,
                color=color, align=PP_ALIGN.RIGHT)


def add_title_block(slide, title, subtitle):
    add_textbox(slide, Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.8),
                title, font=FONT_HEAD, size=34, bold=True, color=BG_NAVY)
    # underline divider
    add_rect(slide, Inches(0.5), Inches(1.25), Inches(12.3), Pt(1.5), GREY_300)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.35), Inches(12.3), Inches(0.5),
                    subtitle, font=FONT_BODY, size=16, italic=True, color=GREY_700)


def add_speaker_notes(slide, notes):
    """Add speaker notes to a slide."""
    slide.notes_slide.notes_text_frame.text = notes


# ---------- slide builders ----------

def build_slide_1(prs):
    """The measurable gap: three stat cards."""
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "The Philippines has a broadband problem",
        "The World Bank measures it on three axes."
    )

    # Three stat cards, equal width
    card_y = Inches(2.1)
    card_h = Inches(3.7)
    card_w = Inches(3.95)
    gap = Inches(0.175)
    start_x = Inches(0.5)

    cards = [
        ("33%", "Access",
         "Household fixed-broadband penetration. Viet Nam 76%. Malaysia 50%."),
        ("11%", "Affordability",
         "Fixed broadband as share of GNI per capita. Highest in ASEAN."),
        ("Low", "Quality",
         "Median fixed speeds below every major ASEAN peer."),
    ]

    for i, (big, label, ctx) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        # card body
        add_rect(slide, x, card_y, card_w, card_h, GREY_100)
        # left accent bar
        add_left_border(slide, x, card_y, card_h, BG_NAVY, thickness_pt=6)

        # big number
        add_textbox(slide, x + Inches(0.3), card_y + Inches(0.3),
                    card_w - Inches(0.6), Inches(1.6),
                    big, font=FONT_HEAD, size=72, bold=True, color=BG_NAVY,
                    line_spacing=1.0)
        # label
        add_textbox(slide, x + Inches(0.3), card_y + Inches(1.95),
                    card_w - Inches(0.6), Inches(0.4),
                    label.upper(), font=FONT_HEAD, size=13, bold=True,
                    color=RISE_BLACK)
        # context
        add_textbox(slide, x + Inches(0.3), card_y + Inches(2.4),
                    card_w - Inches(0.6), card_h - Inches(2.6),
                    ctx, font=FONT_BODY, size=13, color=GREY_700,
                    line_spacing=1.4)

    add_source(slide,
               'Source: World Bank, "Better Internet for All Filipinos" (January 2024). Ookla (2023, 2025).')
    add_rise_mark(slide)

    add_speaker_notes(slide,
        "The Philippines has a broadband problem. The World Bank measures it on "
        "three axes: access, affordability, and quality. Household fixed "
        "penetration, 33 percent. Bandwidth, 11 percent of GNI per capita, the "
        "highest in ASEAN. Speeds lag every major ASEAN peer.")


def build_slide_2(prs):
    """Small ISP margin trap: comparison table."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "Small ISPs carry the expansion",
        "But the margin structure traps them."
    )

    # Table built with shapes for full style control
    table_top = Inches(2.1)
    col_metric_x = Inches(0.5)
    col_metric_w = Inches(6.0)
    col_small_x = Inches(6.6)
    col_small_w = Inches(3.15)
    col_big_x = Inches(9.85)
    col_big_w = Inches(2.95)
    row_h = Inches(0.65)

    # header row
    header_y = table_top
    add_rect(slide, col_metric_x, header_y, SLIDE_W - Inches(1.0), Pt(2),
             GREY_300)  # thin top rule
    add_textbox(slide, col_metric_x, header_y + Inches(0.05),
                col_metric_w, Inches(0.5),
                "FIVE-YEAR AVERAGE, 2017-2022",
                font=FONT_HEAD, size=11, bold=True,
                color=GREY_700, align=PP_ALIGN.LEFT)
    add_textbox(slide, col_small_x, header_y + Inches(0.05),
                col_small_w, Inches(0.5),
                "SMALL ISPs",
                font=FONT_HEAD, size=12, bold=True,
                color=RISE_PINK, align=PP_ALIGN.CENTER)
    add_textbox(slide, col_big_x, header_y + Inches(0.05),
                col_big_w, Inches(0.5),
                "GLOBE / PLDT",
                font=FONT_HEAD, size=12, bold=True,
                color=BG_NAVY, align=PP_ALIGN.CENTER)

    # header rule (bottom of header)
    rule_y = header_y + Inches(0.65)
    add_rect(slide, col_metric_x, rule_y, SLIDE_W - Inches(1.0), Pt(1.5), GREY_500)

    # data rows
    rows = [
        ("Revenue growth", "83%", "29%"),
        ("Net income margin", "4%", "14%"),
        ("CAPEX as % of revenue", "17%", "47%"),
        ("Wholesale bandwidth as % of revenue", "~33%", "—"),
    ]
    for i, (metric, small, big) in enumerate(rows):
        y = rule_y + Inches(0.2) + row_h * i
        # metric label
        add_textbox(slide, col_metric_x, y, col_metric_w, row_h,
                    metric, font=FONT_BODY, size=15, color=RISE_BLACK,
                    anchor=MSO_ANCHOR.MIDDLE)
        # small ISP number
        add_textbox(slide, col_small_x, y, col_small_w, row_h,
                    small, font=FONT_HEAD, size=24, bold=True,
                    color=RISE_PINK, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        # big number
        add_textbox(slide, col_big_x, y, col_big_w, row_h,
                    big, font=FONT_HEAD, size=24, bold=True,
                    color=BG_NAVY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE)
        # row separator
        sep_y = y + row_h
        add_rect(slide, col_metric_x, sep_y, SLIDE_W - Inches(1.0),
                 Pt(0.5), GREY_300)

    # Callout
    callout_y = rule_y + Inches(0.2) + row_h * 4 + Inches(0.3)
    callout_h = Inches(0.85)
    add_rect(slide, col_metric_x, callout_y, SLIDE_W - Inches(1.0),
             callout_h, RGBColor(0xFF, 0xFA, 0xF0))
    add_left_border(slide, col_metric_x, callout_y, callout_h,
                    RGBColor(0xD6, 0x9E, 0x2E), thickness_pt=4)
    add_textbox(slide, col_metric_x + Inches(0.25), callout_y + Inches(0.12),
                SLIDE_W - Inches(1.5), callout_h - Inches(0.25),
                "Many small ISPs began as cable TV operators, now running fiber. "
                "They can't build faster, cheaper, or better because there's "
                "almost no revenue left to reinvest.",
                font=FONT_BODY, size=13, color=RGBColor(0x74, 0x42, 0x10),
                line_spacing=1.3)

    add_source(slide, "Source: World Bank, 2024, based on PSE and SEC financial statements.")
    add_rise_mark(slide)

    add_speaker_notes(slide,
        "Small ISPs carry the expansion. Many began as cable TV operators, now "
        "running fiber. Their revenues grew 83 percent over five years. But "
        "they earn just 4 percent net margin, reinvest only 17 percent of "
        "revenue in CAPEX, and roughly a third of their revenue goes straight "
        "to wholesale bandwidth. That's why they can't build faster, cheaper, "
        "or better.")


def build_slide_3(prs):
    """Capability as lever: causal chain."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "Capability is the lever",
        "Regulatory reform is coming. Capability is controllable today."
    )

    # 5 boxes with arrows between
    box_y = Inches(2.6)
    box_h = Inches(2.2)
    num_boxes = 5
    total_w = SLIDE_W - Inches(1.0)  # 12.3"
    box_w = Inches(2.2)
    gap = (total_w - box_w * num_boxes) / (num_boxes - 1)  # arrow gap
    start_x = Inches(0.5)

    boxes = [
        ("Capability", None),
        ("Peering,\nrouting security,\nIXP participation", None),
        ("Lower\nbandwidth\nspend", None),
        ("Restored\nreinvestment\ncapacity", None),
        ("Access\nAffordability\nQuality", "outcome"),
    ]

    for i, (text, kind) in enumerate(boxes):
        x = start_x + (box_w + gap) * i
        is_outcome = kind == "outcome"
        bg = RGBColor(0xEB, 0xF8, 0xFF) if is_outcome else GREY_100
        top = BG_NAVY if not is_outcome else TEAL
        add_rect(slide, x, box_y, box_w, box_h, bg)
        add_top_border(slide, x, box_y, box_w, top, thickness_pt=4)

        # text
        text_color = TEAL if is_outcome else RISE_BLACK
        weight = True if (i == 0 or is_outcome) else False
        add_textbox(slide, x + Inches(0.1), box_y + Inches(0.2),
                    box_w - Inches(0.2), box_h - Inches(0.4),
                    text, font=FONT_HEAD if weight else FONT_BODY,
                    size=16, bold=weight, color=text_color,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.25)

        # arrow after (skip last)
        if i < num_boxes - 1:
            arrow_x = x + box_w
            arrow_y = box_y + box_h / 2 - Inches(0.15)
            add_textbox(slide, arrow_x, arrow_y, gap, Inches(0.35),
                        "→", font=FONT_HEAD, size=22,
                        color=GREY_500, align=PP_ALIGN.CENTER,
                        anchor=MSO_ANCHOR.MIDDLE)

    # caption below
    add_textbox(slide, Inches(0.5), box_y + box_h + Inches(0.4),
                SLIDE_W - Inches(1.0), Inches(0.4),
                "Reinvestment is what moves the three World Bank indicators.",
                font=FONT_BODY, size=14, italic=True, color=GREY_700,
                align=PP_ALIGN.CENTER)

    add_rise_mark(slide)

    add_speaker_notes(slide,
        "Regulatory reform is coming. In the meantime, capability is the one "
        "lever operators themselves control. Peering, routing security, and "
        "IXP participation reduce the bandwidth bill directly. That restores "
        "reinvestment capacity. And reinvestment is what moves the three "
        "World Bank indicators.")


def build_slide_4(prs):
    """IXperience Pro: curriculum + TTT + photo placeholder."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "IXperience Pro",
        "A structured cohort program for provincial Philippine engineers, co-developed with APNIC Academy."
    )

    # 4 module cards in a row
    card_y = Inches(2.1)
    card_h = Inches(1.6)
    num_cards = 4
    total_w = SLIDE_W - Inches(1.0)
    card_gap = Inches(0.15)
    card_w = (total_w - card_gap * (num_cards - 1)) / num_cards
    start_x = Inches(0.5)

    modules = [
        ("MODULE 1", "BGP fundamentals"),
        ("MODULE 2", "Peering and interconnection"),
        ("MODULE 3", "Routing security (RPKI / MANRS)"),
        ("MODULE 4", "IXP operations"),
    ]
    for i, (num, name) in enumerate(modules):
        x = start_x + (card_w + card_gap) * i
        add_rect(slide, x, card_y, card_w, card_h, GREY_100)
        add_top_border(slide, x, card_y, card_w, BG_NAVY, thickness_pt=4)
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.25),
                    card_w - Inches(0.4), Inches(0.3),
                    num, font=FONT_HEAD, size=10, bold=True,
                    color=GREY_500)
        add_textbox(slide, x + Inches(0.2), card_y + Inches(0.55),
                    card_w - Inches(0.4), card_h - Inches(0.7),
                    name, font=FONT_HEAD, size=18, bold=True,
                    color=BG_NAVY, line_spacing=1.25)

    # TTT strip
    ttt_y = card_y + card_h + Inches(0.3)
    ttt_h = Inches(1.0)
    add_rect(slide, Inches(0.5), ttt_y, SLIDE_W - Inches(1.0), ttt_h,
             RGBColor(0xFD, 0xEB, 0xEE))  # soft pink tint
    add_left_border(slide, Inches(0.5), ttt_y, ttt_h, RISE_PINK, thickness_pt=5)
    add_textbox(slide, Inches(0.75), ttt_y + Inches(0.15),
                SLIDE_W - Inches(1.5), ttt_h - Inches(0.3),
                "4 cohorts × 15 engineers = 60 certified. Train-the-trainer "
                "track certifies 4 to 6 Filipino facilitators, so delivery "
                "becomes independent of international instructors. All "
                "materials publish under Creative Commons.",
                font=FONT_BODY, size=14, color=RISE_BLACK, line_spacing=1.35)

    # Photo placeholder
    ph_y = ttt_y + ttt_h + Inches(0.25)
    ph_h = Inches(1.2)
    ph_rect = add_rect(slide, Inches(0.5), ph_y,
                       SLIDE_W - Inches(1.0), ph_h, GREY_100)
    ph_rect.line.color.rgb = GREY_300
    ph_rect.line.width = Pt(1.5)
    ph_rect.line.dash_style = 7  # dash
    add_textbox(slide, Inches(0.5), ph_y, SLIDE_W - Inches(1.0), ph_h,
                "[ PHOTO PLACEHOLDER: past IXperience workshop, participants in learning mode ]",
                font=FONT_BODY, size=12, italic=True, color=GREY_500,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rise_mark(slide)

    add_speaker_notes(slide,
        "IXperience Pro is a structured cohort program for provincial "
        "Philippine engineers. Four modules: BGP, peering, routing security, "
        "and IXP operations, co-developed with APNIC Academy. Four cohorts of "
        "fifteen, sixty engineers from at least twenty ISPs. And a "
        "train-the-trainer track that certifies four to six Filipino "
        "facilitators, so delivery becomes independent of international "
        "instructors.")


def build_slide_5(prs):
    """Not an experiment: track record + credentials + photo."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, WHITE)

    add_title_block(
        slide,
        "Not an experiment",
        "Scaling something already working."
    )

    # Four stat cards
    stat_y = Inches(2.1)
    stat_h = Inches(1.4)
    num_stats = 4
    total_w = SLIDE_W - Inches(1.0)
    stat_gap = Inches(0.15)
    stat_w = (total_w - stat_gap * (num_stats - 1)) / num_stats
    start_x = Inches(0.5)

    stats = [
        ("4", "Prior events"),
        ("155", "Companies"),
        ("251", "Engineers trained"),
        ("71%", "Provincial"),
    ]
    for i, (num, label) in enumerate(stats):
        x = start_x + (stat_w + stat_gap) * i
        add_rect(slide, x, stat_y, stat_w, stat_h, GREY_100)
        add_top_border(slide, x, stat_y, stat_w, BG_NAVY, thickness_pt=3)
        add_textbox(slide, x, stat_y + Inches(0.2),
                    stat_w, Inches(0.7),
                    num, font=FONT_HEAD, size=40, bold=True,
                    color=BG_NAVY, align=PP_ALIGN.CENTER,
                    line_spacing=1.0)
        add_textbox(slide, x, stat_y + Inches(0.95),
                    stat_w, Inches(0.4),
                    label.upper(), font=FONT_HEAD, size=11, bold=True,
                    color=GREY_700, align=PP_ALIGN.CENTER)

    # Three credential pills
    pill_y_start = stat_y + stat_h + Inches(0.3)
    pill_h = Inches(0.55)
    pill_gap = Inches(0.12)
    pills = [
        ("Curriculum runs on APNIC Academy labs and courseware.",
         " With APNIC and PhNOG as co-delivery partners."),
        ("2023 ISIF grant delivered:",
         " Project School Internet, 50+ schools connected in Cebu."),
        ("Standing MOA with FICTAP:",
         " RISE serves as the community's designated capability partner."),
    ]
    for i, (bold_part, rest) in enumerate(pills):
        y = pill_y_start + (pill_h + pill_gap) * i
        add_rect(slide, Inches(0.5), y, SLIDE_W - Inches(1.0), pill_h,
                 RGBColor(0xED, 0xF2, 0xF7))
        add_left_border(slide, Inches(0.5), y, pill_h, BG_NAVY, thickness_pt=3)

        # build pill text as two runs (bold + regular)
        box = slide.shapes.add_textbox(Inches(0.75), y + Inches(0.1),
                                        SLIDE_W - Inches(1.5), pill_h - Inches(0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        r1 = p.add_run()
        r1.text = bold_part
        r1.font.name = FONT_HEAD
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = BG_NAVY
        r2 = p.add_run()
        r2.text = rest
        r2.font.name = FONT_BODY
        r2.font.size = Pt(13)
        r2.font.color.rgb = RISE_BLACK

    # Photo placeholder (smaller, at bottom)
    ph_y = pill_y_start + (pill_h + pill_gap) * 3 + Inches(0.1)
    ph_h = Inches(0.7)
    ph_rect = add_rect(slide, Inches(0.5), ph_y,
                       SLIDE_W - Inches(1.0), ph_h, GREY_100)
    ph_rect.line.color.rgb = GREY_300
    ph_rect.line.width = Pt(1.5)
    ph_rect.line.dash_style = 7
    add_textbox(slide, Inches(0.5), ph_y, SLIDE_W - Inches(1.0), ph_h,
                "[ PHOTO PLACEHOLDER: past IXperience participant mosaic ]",
                font=FONT_BODY, size=11, italic=True, color=GREY_500,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_rise_mark(slide)

    add_speaker_notes(slide,
        "This isn't an experiment. With APNIC and PhNOG, we've trained 251 "
        "engineers from 155 companies, seventy-one percent provincial. "
        "Curriculum runs on APNIC Academy labs. Our 2023 ISIF grant, Project "
        "School Internet, connected over 50 schools in Cebu. And under "
        "standing MOA with FICTAP, the fiber ISP trade association, RISE "
        "serves as the community's designated capability partner.")


def build_slide_6(prs):
    """Closing: the loop. Dark RISE navy background."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    set_slide_bg(slide, BG_DARK)

    # Title
    add_textbox(slide, Inches(0.75), Inches(0.9),
                SLIDE_W - Inches(1.5), Inches(1.0),
                "The loop",
                font=FONT_HEAD, size=44, bold=True, color=WHITE)
    # thin divider
    add_rect(slide, Inches(0.75), Inches(1.85),
             Inches(2.5), Pt(2), RISE_PINK)

    # Four bullets with pink arrow glyph
    bullets = [
        "Filipino engineers trained by Filipino engineers",
        "Lower bandwidth costs",
        "Restored reinvestment capacity",
        "Measurable movement on the indicators the World Bank tracks",
    ]
    bullet_start_y = Inches(2.5)
    bullet_h = Inches(0.65)
    for i, text in enumerate(bullets):
        y = bullet_start_y + bullet_h * i
        # arrow glyph
        add_textbox(slide, Inches(0.75), y, Inches(0.6), bullet_h,
                    "→", font=FONT_HEAD, size=24, bold=True,
                    color=RISE_PINK, anchor=MSO_ANCHOR.MIDDLE)
        # bullet text
        add_textbox(slide, Inches(1.35), y,
                    SLIDE_W - Inches(2.1), bullet_h,
                    text, font=FONT_BODY, size=22, color=WHITE,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.2)

    # Partner row at bottom, separator line above
    partner_y = Inches(6.6)
    add_rect(slide, Inches(0.75), partner_y, SLIDE_W - Inches(1.5), Pt(1),
             RGBColor(0x2D, 0x37, 0x48))
    add_textbox(slide, Inches(0.75), partner_y + Inches(0.2),
                SLIDE_W - Inches(1.5), Inches(0.4),
                "RISE      APNIC ACADEMY      PhNOG      FICTAP",
                font=FONT_HEAD, size=13, bold=True,
                color=GREY_500, align=PP_ALIGN.LEFT)

    add_speaker_notes(slide,
        "Filipino engineers trained by Filipino engineers. Lower bandwidth "
        "costs. Restored reinvestment capacity. Measurable movement on the "
        "indicators the World Bank tracks.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide_1(prs)
    build_slide_2(prs)
    build_slide_3(prs)
    build_slide_4(prs)
    build_slide_5(prs)
    build_slide_6(prs)

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}")
    print(f"  {len(prs.slides)} slides, native editable content")


if __name__ == "__main__":
    main()
